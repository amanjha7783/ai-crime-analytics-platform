from __future__ import annotations

from pathlib import Path


SLIDES = [
    ("AI-Driven Crime Analytics", "Visualization Platform for Karnataka State Police Datathon 2026"),
    ("Problem", "Fragmented police records limit pattern detection, offender tracking, network analysis, hotspot prediction, and proactive policing."),
    ("Solution", "A full intelligence platform with ETL, PostGIS, machine learning, explainability, FastAPI, and a Next.js police command dashboard."),
    ("Data Pipeline", "Collect, validate, clean, impute, engineer features, analyze, store, predict, explain, and visualize."),
    ("ML Intelligence", "Classification, hotspot prediction, repeat offender detection, risk scoring, trend forecasting, anomaly detection, and explainability."),
    ("Dashboard", "KPIs, maps, heatmaps, trends, criminal network graph, risk scores, filters, and reports."),
    ("Architecture", "Police data flows through ETL to PostGIS and ML services, then through FastAPI to the intelligence portal."),
    ("Impact", "Prioritize patrols, detect repeat activity, compare districts, investigate networks, and produce faster intelligence briefs."),
    ("Deployment", "Docker locally, Vercel for frontend, Render or Railway for backend, PostgreSQL/PostGIS, and Redis."),
    ("Future Scope", "Official secure feeds, advanced models, SHAP, streaming alerts, investigator workflows, and model monitoring."),
]


def build_presentation(output_path: Path) -> None:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError as exc:
        raise SystemExit("Install python-pptx with: python -m pip install python-pptx") from exc

    presentation = Presentation()
    for title, body in SLIDES:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.7), Inches(8.6), Inches(1.2))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(34)
        title_frame.paragraphs[0].font.bold = True
        body_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.1), Inches(8.4), Inches(3.8))
        body_frame = body_box.text_frame
        body_frame.text = body
        body_frame.paragraphs[0].font.size = Pt(22)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output_path)


if __name__ == "__main__":
    build_presentation(Path("docs/datathon-presentation.pptx"))
